from pptx import Presentation
from pptx.util import Inches, Pt
import os

# Пътища
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
BASE_DIR = os.path.dirname(SCRIPT_DIR)
PRESENTATION_DIR = os.path.join(BASE_DIR, 'presentation')

def create_presentation():
    prs = Presentation()
    
    # Slide 0: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Академичен анализ на Дървовидните Алгоритми"
    slide.placeholders[1].text = "От CART до Gradient Boosting\nОснови на изкуствения интелект"
    
    # Slide 1: Въведение и парадигма
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "1. Въведение и Парадигма"
    tf = slide.placeholders[1].text_frame
    tf.text = "Дърветата на решенията са базирани на принципа 'Разделяй и владей' (Divide and Conquer)."
    p = tf.add_paragraph()
    p.text = "Алгоритъмът разчленява пространството на хипер-правоъгълници."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "CART (Classification and Regression Trees) обединява двата основни подхода:"
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Класификация: Предсказва дискретна категория."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Регресия: Предсказва конкретна числова стойност."
    p.level = 2

    # Slide 2: Анатомия на дървото
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "2. Анатомия на Дървото"
    tf = slide.placeholders[1].text_frame
    tf.text = "Структурни елементи на модела:"
    p = tf.add_paragraph()
    p.text = "Корен (Root Node): Началната точка с цялата база данни."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Вътрешни възли (Internal Nodes): Задават логически въпроси (тестват атрибути)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Клони (Branches): Представляват отговорите (най-често бинарни Да/Не)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Листа (Leaf Nodes): Крайните точки, съдържащи прогнозата на модела."
    p.level = 1

    # Slide 3: Класически алгоритми
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3. Класически алгоритми"
    tf = slide.placeholders[1].text_frame
    tf.text = "Еволюция на архитектурите:"
    p = tf.add_paragraph()
    p.text = "ID3 (Iterative Dichotomiser 3): Използва Information Gain; само за категорийни данни."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "C4.5: Въвежда числови стойности, подрязване (Pruning) и Gain Ratio."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "CART: Съвременният стандарт. Генерира бинарни дървета; използва Gini за класификация и MSE за регресия."
    p.level = 1

    # Slide 4: Математически фундамент - Класификация
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "4. Математика: Класификация"
    tf = slide.placeholders[1].text_frame
    tf.text = "Цел: Постигане на хомогенност (чистота) в нодовете."
    p = tf.add_paragraph()
    p.text = "Gini Impurity: Минимизира вероятността за грешна класификация."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Формула: Gini = 1 - Σ (p_i)²"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Ентропия (Entropy): Измерва хаоса в данните (по Шанън)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Формула: Entropy = -Σ p_i * log2(p_i)"
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Information Gain: Намалението на ентропията след разделяне."
    p.level = 2

    # Slide 5: Математически фундамент - Регресия
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "5. Математика: Регресия"
    tf = slide.placeholders[1].text_frame
    tf.text = "Цел: Минимизиране на дисперсията и грешката на предсказанието (средна стойност)."
    p = tf.add_paragraph()
    p.text = "MSE (Mean Squared Error): Средноквадратична грешка."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Силно наказва големите отклонения."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "MAE (Mean Absolute Error): Средна абсолютна грешка."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "По-устойчива на екстремни стойности (outliers)."
    p.level = 2

    # Slide 6: Преобучаване и Регуларизация
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "6. Преобучаване и Регуларизация"
    tf = slide.placeholders[1].text_frame
    tf.text = "Проблем: Дърветата са склонни да 'наизустяват' шума в данните (Overfitting)."
    p = tf.add_paragraph()
    p.text = "Пре-подрязване (Pre-pruning): Спиране на растежа."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Max Depth, Min Samples Split, Min Samples Leaf."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Пост-подрязване (Post-pruning): Редуциране след пълно израстване."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Cost Complexity Pruning."
    p.level = 2

    # Slide 7: Ансамблово обучение - Bagging
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "7. Ансамбли: Random Forest (Bagging)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Комбиниране на 'слаби' модели за създаване на 'силен'."
    p = tf.add_paragraph()
    p.text = "Концепция: Паралелно обучение на стотици дървета."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Всяко дърво се обучава на случайна извадка от данните (Bootstrap)."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Използва случайно подмножество от признаците (Де-корелация)."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Крайният резултат е мажоритарно гласуване или осредняване."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Решава проблема с нестабилността на единичното дърво."
    p.level = 2

    # Slide 8: Ансамблово обучение - Boosting
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "8. Ансамбли: XGBoost (Boosting)"
    tf = slide.placeholders[1].text_frame
    tf.text = "Фокус върху оптимизацията и прецизността."
    p = tf.add_paragraph()
    p.text = "Концепция: Последователно обучение на дървета."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Всяко ново дърво коригира грешките (остатъците) на предходното."
    p.level = 2
    p = tf.add_paragraph()
    p.text = "Gradient Boosting: Използва Gradient Descent за минимизиране на загубната функция."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Изключително висока точност; доминира състезанията за таблични данни."
    p.level = 2

    # Slide 9: Приложения в индустрията
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "9. Индустриални приложения"
    tf = slide.placeholders[1].text_frame
    tf.text = "Дървовидните модели управляват критични глобални системи:"
    p = tf.add_paragraph()
    p.text = "Netflix / Amazon: Ранкинг в препоръчителните системи (XGBoost)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "PayPal: Откриване на измами (Fraud Detection) в реално време (Random Forest)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "NASA: Класификация на небесни тела; астрономите се нуждаят от интерпретируемостта на дърветата."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Uber: Предсказване на точно време на пристигане (ETA)."
    p.level = 1

    # Slide 10: Заключение
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "10. Заключение"
    tf = slide.placeholders[1].text_frame
    tf.text = "Дървовидните алгоритми предлагат уникален баланс:"
    p = tf.add_paragraph()
    p.text = "Simple Trees: Максимална прозрачност и лесен одит (Бяла кутия)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Random Forest: Изключителна стабилност и устойчивост на шум."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "XGBoost: Ненадмината изчислителна мощ и математическа прецизност."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Те остават незаменим инструмент в арсенала на изкуствения интелект."
    p.level = 1

    if not os.path.exists(PRESENTATION_DIR):
        os.makedirs(PRESENTATION_DIR)
        
    save_path = os.path.join(PRESENTATION_DIR, 'Decision_Trees_Theory_Presentation.pptx')
    try:
        prs.save(save_path)
        print(f"Теоретичната презентация е успешно генерирана: {save_path}")
    except Exception as e:
        print(f"Грешка при запис: {e}")

if __name__ == '__main__':
    create_presentation()
